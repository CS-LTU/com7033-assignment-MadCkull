import collections 
import collections.abc

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    Presentation = None


# Define Slide Content Structure
SLIDES = [
    {
        "title": "StrokeVision: Decision Support System",
        "subtitle": "A Secure, AI-Driven Clinical Web Application",
        "content": ["Presenter: [Your Name]", "Duration: 15 Minutes"],
        "notes": "Welcome everyone to the StrokeVision presentation. Today I will demonstrate how we integrated deep learning with secure clinical data management."
    },
    {
        "title": "Project Overview & Objectives",
        "content": [
            "Goal: Early detection of stroke risk using ML on patient vitals.",
            "Key Features:",
            "  - Role-Based Access Control (RBAC)",
            "  - Secure CRUD for Patient Records",
            "  - Real-time Risk Prediction (Keras/TensorFlow)",
            "Target Users: Doctors, Nurses, Admin"
        ],
        "notes": "StrokeVision empowers healthcare professionals with instant risk assessments while maintaining strict data privacy."
    },
    {
        "title": "System Design & Architecture (SSDLC)",
        "content": [
            "Frontend: Single Page Application (Vanilla JS)",
            "Backend: Flask REST API (Secure Gateway)",
            "Databases (Separation of Concerns):",
            "  - SQLite: User Auth & Credentials",
            "  - MongoDB: Patient Health Records",
            "SSDLC: Security integrated at Design, Implementation, and Testing."
        ],
        "notes": "We utilized a decoupled architecture. Separating auth data (SQLite) from patient data (MongoDB) minimizes the blast radius of any potential breach."
    },
    {
        "title": "Security Features Overview",
        "content": [
            "Authentication: Hybrid Session (Flask-Login) + JWT",
            "Data Protection: Bcrypt Password Hashing (Work Factor 12)",
            "Input Validation: Server-Side (Flask-WTF) & Sanitization",
            "Network Security: CSRF Protection & Secure Headers",
            "Session Management: HttpOnly, Secure Cookies"
        ],
        "notes": "Security is the backbone. We use Bcrypt for passwords and enforce strict server-side validation to prevent Injection attacks."
    },
    {
        "title": "Database Security & Privacy",
        "content": [
            "Database Separation Strategy:",
            "  - Users.db (SQLite) vs StrokeVision_DB (MongoDB)",
            "Access Control: Least Privilege Connection Strings",
            "Encryption Strategy:",
            "  - Implemented: Password Hashing",
            "  - Planned: Field-Level AES Encryption for Patient Data"
        ],
        "notes": "Our dual-database entry strategy ensures granular access control. We openly acknowledge that field-level encryption is the next planned feature."
    },
    {
        "title": "Ethical & Professional Development",
        "content": [
            "Data Privacy (GDPR/UK DPA):",
            "  - Data Minimization (Only vital fields collected)",
            "  - Accountability (Audit properties on records)",
            "Professional Practices:",
            "  - Version Control (Git/GitHub)",
            "  - Automated Testing (pytest)",
            "  - Comprehensive Documentation"
        ],
        "notes": "We adhered to GDPR principles like data minimization. The project was managed professionally using Git and CI/CD best practices."
    },
    {
        "title": "Demonstration Overview (Live)",
        "content": [
            "Scenario: Full Lifecycle Demo",
            "1. Secure Registration & Login",
            "2. Role-Based Access Control (Admin blocks)",
            "3. CRUD Operations & Validation",
            "4. Security Defense (XSS/CSRF)",
            "5. Audit Logs Review"
        ],
        "notes": "I will now proceed to the live demo, showing the system in action."
    },
    {
        "title": "Demo: Authentication & RBAC",
        "content": [
            "Objective: Verify Access Controls",
            "Evidence:",
            "  - Secure Redirects after Login",
            "  - 403 Forbidden on Unauthorized Access",
            "  - HttpOnly Cookie inspection"
        ],
        "notes": "Notice how the system securely handles sessions and blocks unauthorized role access."
    },
    {
        "title": "Demo: Patient Management & Validation",
        "content": [
            "Objective: Verify Data Integrity",
            "Evidence:",
            "  - Inputs Rejected (Invalid Age/Formats)",
            "  - XSS Payloads Sanitized",
            "  - Successful MongoDB Persistence"
        ],
        "notes": "The input validation layer is our first line of defense against injection attacks."
    },
    {
        "title": "Demo: Testing & Code Quality",
        "content": [
            "Objective: Verify Reliability",
            "Evidence:",
            "  - Passing `pytest` Suite",
            "  - Structured GitHub Commit History",
            "  - Clean Code Architecture"
        ],
        "notes": "Our test suite ensures that security logic remains intact even as features are added."
    },
    {
        "title": "OWASP Alignment",
        "content": [
            "A01: Broken Access Control -> RBAC Implementation",
            "A02: Cryptographic Failures -> Bcrypt & HTTPS",
            "A03: Injection -> ORM & WTForms",
            "A04: Insecure Design -> Threat Modeling"
        ],
        "notes": "We explicitly mapped our defenses against the OWASP Top 10 vulnerabilities."
    },
    {
        "title": "Reflections & Critique",
        "content": [
            "Strengths: Robust Auth, Clean UI, Strong Architecture.",
            "Limitations:",
            "  - Missing Field-Level Encryption",
            "  - Lack of MFA",
            "Challenges: Managing Dual-Database Consistency"
        ],
        "notes": "The project is strong, but we recognize the need for field-level encryption in the next sprint."
    },
    {
        "title": "Future Roadmap",
        "content": [
            "Immediate: Implement AES Encryption for Patient PII",
            "Short-Term: Add Multi-Factor Authentication",
            "Long-Term: Docker Containerization",
            "Testing: End-to-End Browser Testing"
        ],
        "notes": "Our roadmap is clear: enhance data-at-rest security and improve deployment portability."
    },
    {
        "title": "Conclusion",
        "content": [
            "Summary: Professional, Secure, Compliant.",
            "Key Takeaway: Security is integral to the design.",
            "Thank You / Q&A"
        ],
        "notes": "Thank you for listening. I am happy to take any questions."
    }
]

def create_presentation():
    prs = Presentation()

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = SLIDES[0]["title"]
    subtitle.text = SLIDES[0]["subtitle"] + "\n" + "\n".join(SLIDES[0]["content"])
    slide.notes_slide.notes_text_frame.text = SLIDES[0]["notes"]

    # Content Slides
    bullet_slide_layout = prs.slide_layouts[1]

    for slide_data in SLIDES[1:]:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = slide_data["title"]
        tf = body_shape.text_frame
        tf.text = slide_data["content"][0] if slide_data["content"] else ""

        for item in slide_data["content"][1:]:
            p = tf.add_paragraph()
            p.text = item
            if item.strip().startswith("-"):
                p.level = 1
            if item.strip().startswith("  -"):
                p.level = 2

        if "notes" in slide_data:
            slide.notes_slide.notes_text_frame.text = slide_data["notes"]

    output_path = "StrokeVision_Presentation.pptx"
    prs.save(output_path)
    print(f"Successfully created presentation: {output_path}")

if __name__ == "__main__":
    print("Generating StrokeVision Presentation...")
    if Presentation is None:
        print("Error: python-pptx is not installed.")
        print("Please run: pip install python-pptx")
    else:
        try:
            create_presentation()
        except Exception as e:
            print(f"An error occurred: {e}")
